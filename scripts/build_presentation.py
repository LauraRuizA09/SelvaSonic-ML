#!/usr/bin/env python3
"""
scripts/build_presentation.py
Genera: results/presentacion/SelvaSonic_ML_presentacion_final.pptx  (23 slides)
Ejecutar con el venv del proyecto:
    .\venv\Scripts\python.exe scripts\build_presentation.py
"""
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

OUT_DIR = ROOT / "results" / "presentacion"
OUT_DIR.mkdir(parents=True, exist_ok=True)

# ── Paleta SelvaSonic-ML ──────────────────────────────────────────────────────
BG   = "00100C"; BG2  = "082C2C"; PNL  = "0B2B27"; PNL2 = "0F302B"
CYN  = "21D4D4"; GRN  = "33D17A"; GLD  = "F4D35E"
ORG  = "FF9F1C"; RED  = "FF5C7A"
TP   = "EFFCF5"; TS   = "B5D3C3"; TM   = "7EAA95"
FONT = "Aptos Display"

# ── Datos (de summary.json / resumen_final.json / history.json) ───────────────
# ECE v2 = 0.1452  > v1 = 0.0389: hallazgo real (class weights deterioran ECE global)
METRICS = {
    "Baseline": dict(
        val_acc="0.7419", test_acc="0.6322", macro_f1="0.509",
        macro_p="0.525",  macro_r="0.536",  ece="0.050",
        params="422,635", epochs="27 / 30",  gap="0.1097",
        train="~286 min",  ap_raras="0.456",
    ),
    "Attention v1": dict(
        val_acc="0.7829", test_acc="0.6873", macro_f1="0.599",
        macro_p="0.606",  macro_r="0.618",  ece="0.039",
        params="714,987", epochs="17 / 30",  gap="0.0956",
        train="~181 min",  ap_raras="0.573",
    ),
    "Attention v2": dict(
        val_acc="0.7653", test_acc="0.7036", macro_f1="0.625",
        macro_p="0.638",  macro_r="0.662",  ece="0.145",
        params="714,987", epochs="25 / 30",  gap="0.0617",
        train="~270 min",  ap_raras="0.664",
    ),
}

SPECIES = [
    ("Trogon viridis",         "Trogonidae",     79,  2784, "0.4693"),
    ("Crypturellus cinereus",  "Tinamidae",      48,  2656, "0.4919"),
    ("Lipaugus vociferans",    "Cotingidae",     36,  1650, "0.7918"),
    ("Ramphastos tucanus",     "Ramphastidae",   31,  1492, "0.8757"),
    ("Crypturellus undulatus", "Tinamidae",      29,   738, "1.7704"),
    ("Celeus grammicus",       "Picidae",        28,   856, "1.5263"),
    ("Glaucidium brasilianum", "Strigidae",      22,   560, "2.3331"),
    ("Chordeiles pusillus",    "Caprimulgidae",  21,   240, "5.4439"),
    ("Frederickena fulva",     "Thamnophilidae", 20,   854, "1.5299"),
    ("Rupornis magnirostris",  "Accipitridae",   20,   302, "4.3263"),
]

IMGS = {
    "dist_clases":  OUT_DIR / "dist_clases.png",
    "mel_ejemplo":  OUT_DIR / "mel_ejemplo.png",
    "curvas_v2":    ROOT / "results/runs/attention_S4_v2_20260602_1332/curvas_entrenamiento.png",
    "confusion_3":  ROOT / "results/comparativa/confusion_matrices.png",
    "calibracion":  ROOT / "results/comparativa/calibracion.png",
    "embeddings":   ROOT / "results/embeddings/visualizacion_embeddings.png",
    "gradcam_cine": ROOT / "results/gradcam/correctas_Crypturellus_cinereus.png",
    "gradcam_q":    ROOT / "results/gradcam/analisis_cuantitativo.png",
    "ext_dist":     ROOT / "results/test_externos/distribuciones_por_audio.png",
    "f1_clase":     ROOT / "results/comparativa/f1_por_clase.png",
}

JUNGLE_BG_COMPOSED = OUT_DIR / "jungle_bg.png"
JUNGLE_SRC         = ROOT / "assets/selvasonic_banner.png"

IMGS["arch_baseline"]  = OUT_DIR / "arch_baseline.png"
IMGS["arch_attention"] = OUT_DIR / "arch_attention.png"

# ── Generar imágenes que no existen en disco ──────────────────────────────────

def gen_dist_clases():
    import torch, matplotlib; matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    cw  = torch.load(ROOT / "data/class_weights.pt", weights_only=False)
    ct  = cw["class_counts"]; lm = cw["label_map"]
    i2n = {v: k for k, v in lm.items()}
    cc  = {i2n[i]: int(ct[i]) for i in range(len(ct))}
    scc = sorted(cc.items(), key=lambda x: x[1], reverse=True)
    cls = [c for c, _ in scc]; cnt = [n for _, n in scc]
    col = ["#7EAA95" if c == "no_ave" else "#21D4D4" for c in cls]
    fig, ax = plt.subplots(figsize=(10, 5.2), facecolor="#00100C")
    ax.set_facecolor("#00100C")
    ax.barh(range(len(cls)), cnt, color=col, alpha=0.9, edgecolor="none")
    ax.set_yticks(range(len(cls)))
    ax.set_yticklabels([c.replace("_", " ") for c in cls], color="#EFFCF5", fontsize=9)
    ax.invert_yaxis()
    ax.set_xlabel("Clips de entrenamiento", color="#B5D3C3", fontsize=9)
    ax.tick_params(colors="#B5D3C3")
    for sp in ["top", "right"]: ax.spines[sp].set_visible(False)
    for sp in ["bottom", "left"]:
        ax.spines[sp].set_color("#7EAA95"); ax.spines[sp].set_alpha(0.4)
    ax.grid(axis="x", color="#7EAA95", alpha=0.15, linewidth=0.5)
    mx = max(cnt)
    for i, c in enumerate(cnt):
        ax.text(c + mx * 0.01, i, f"{c:,}", va="center", color="#EFFCF5", fontsize=8)
    plt.tight_layout(pad=0.4)
    plt.savefig(IMGS["dist_clases"], dpi=200, facecolor="#00100C", bbox_inches="tight")
    plt.close(); print("  [gen] dist_clases.png")


def gen_arch_diagrams():
    """
    Genera arch_baseline.png y arch_attention.png.
    Estilo: stacked feature-map pancakes (como imagen CNN referencia) + dark SelvaSonic theme.
    """
    import matplotlib; matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.patches import FancyBboxPatch

    _D="#00100C"; _P="#0B2B27"; _B="#082C2C"
    _C="#21D4D4"; _G="#33D17A"; _L="#F4D35E"
    _O="#FF9F1C"; _R="#FF5C7A"; _W="#EFFCF5"
    _S="#B5D3C3"; _M="#7EAA95"

    def _txt(ax, cx, cy, lines):
        n = len(lines); step = 0.32
        top = cy + (n - 1) / 2.0 * step
        for i, (t, sz, col, bd) in enumerate(lines):
            ax.text(cx, top - i*step, t, ha="center", va="center",
                    fontsize=sz, color=col,
                    fontweight="bold" if bd else "normal", zorder=9)

    def fmap(ax, x, yc, w, h, n, ec, lines):
        """Stacked pancakes — simula profundidad de canales."""
        dx, dy = 0.068, 0.036
        for i in range(n-1, -1, -1):
            ax.add_patch(FancyBboxPatch(
                (x+i*dx, yc-h/2-i*dy), w, h,
                boxstyle="square,pad=0", facecolor=_P, edgecolor=ec,
                lw=1.6 if i==0 else 0.5,
                alpha=1.0 if i==0 else 0.40, zorder=4-i))
        _txt(ax, x+w/2, yc, lines)

    def sbox(ax, x, yc, w, h, ec, lines, fc=None):
        ax.add_patch(FancyBboxPatch(
            (x, yc-h/2), w, h, boxstyle="round,pad=0.07",
            facecolor=fc or _B, edgecolor=ec, lw=2.2, zorder=3))
        _txt(ax, x+w/2, yc, lines)

    def arr(ax, x0, x1, y, col=_M):
        ax.annotate("", xy=(x1,y), xytext=(x0,y),
            arrowprops=dict(arrowstyle="-|>", color=col, lw=1.6, mutation_scale=12), zorder=8)

    def brace(ax, x0, x1, y, label, col=_M):
        ax.annotate("", xy=(x1,y), xytext=(x0,y),
            arrowprops=dict(arrowstyle="<->", color=col, lw=0.9, mutation_scale=8), zorder=8)
        ax.text((x0+x1)/2, y-0.22, label, ha="center",
                fontsize=8, color=col, style="italic", zorder=9)

    SPECIES = ["Trogon vir.", "C. cinereus", "Lipaugus", "Ramphastos",
               "C. undulatus", "Celeus gram.", "Glaucidium", "Chordeiles",
               "Frederickena", "Rupornis mag.", "no_ave"]
    Y=2.05; HF=2.65; HB=1.45

    # ════════════════════════════ BASELINE ════════════════════════════════════
    fig, ax = plt.subplots(figsize=(14.8, 4.1), facecolor=_D)
    ax.set_facecolor(_D); ax.set_xlim(0,14.8); ax.set_ylim(0,4.1); ax.axis("off")

    ax.text(7.1, 3.93, "Baseline CNN — SelvaSonic-ML    422,635 parámetros",
            ha="center", fontsize=13, color=_W, fontweight="bold")

    sbox(ax, 0.12, Y, 1.2, HB, _C,
         [("Input",10,_C,True),("[B,1,128,216]",8,_S,False),("Mel-espectro",7.5,_M,False)])
    arr(ax, 1.32,1.55,Y)

    _cb = [
        (1.55, 2, _G, [("Conv 1",9,_G,True),("32 ch",8,_S,False),("BN·ReLU",7,_M,False),("MaxPool 2×2",7,_M,False)]),
        (3.35, 3, _G, [("Conv 2",9,_G,True),("64 ch",8,_S,False),("BN·ReLU",7,_M,False),("MaxPool 2×2",7,_M,False)]),
        (5.25, 3, _G, [("Conv 3",9,_G,True),("128 ch",8,_S,False),("BN·ReLU",7,_M,False),("MaxPool 2×2",7,_M,False)]),
        (7.15, 4, _G, [("Conv 4",9,_G,True),("256 ch",8,_S,False),("BN·ReLU",7,_M,False),("MaxPool 2×2",7,_M,False)]),
    ]
    for i,(x,ns,col,lines) in enumerate(_cb):
        fmap(ax, x, Y, 1.05, HF, ns, col, lines)
        xe = x+1.05+(ns-1)*0.068
        if i<3: arr(ax, xe+0.04, xe+0.19, Y)

    brace(ax, 1.55, 8.55, 0.42, "Feature Extraction  (4 × Conv2d + BatchNorm + ReLU + MaxPool 2×2)")

    xe4 = 7.15+1.05+3*0.068  # ≈ 8.404
    arr(ax, 8.42, 8.68, Y)
    sbox(ax, 8.68, Y, 1.1, 1.25, _L,
         [("GAP",10.5,_L,True),("[B, 256]",8,_S,False),("Global Avg Pool",7,_M,False)])
    arr(ax, 9.78, 10.02, Y)

    sbox(ax, 10.02, Y, 1.55, HB, _O,
         [("FC  256→128",9.5,_O,True),("ReLU",8,_S,False),("Dropout  0.3",7.5,_M,False)])
    arr(ax, 11.57, 11.82, Y)

    sbox(ax, 11.82, Y, 1.45, 1.65, _R,
         [("FC  128→11",9.5,_R,True),("LogSoftmax",8,_S,False),("11 clases",7.5,_M,False)])
    arr(ax, 13.27, 13.48, Y)

    oh = HF/len(SPECIES)
    for k,sp in enumerate(SPECIES):
        yo = Y-HF/2+k*oh
        fc2 = "#112F27" if sp!="no_ave" else "#0D2620"
        ec2 = _M if sp!="no_ave" else _C
        lw2 = 0.7 if sp!="no_ave" else 1.6
        ax.add_patch(FancyBboxPatch((13.48,yo+0.01),1.28,oh-0.02,
            boxstyle="square,pad=0",facecolor=fc2,edgecolor=ec2,lw=lw2,zorder=3))
        ax.text(14.12,yo+oh/2,sp,ha="center",va="center",fontsize=5.8,color=_W,zorder=9)

    brace(ax, 8.68, 14.78, 0.42, "Clasificación")

    plt.tight_layout(pad=0.18)
    plt.savefig(str(OUT_DIR/"arch_baseline.png"), dpi=180, facecolor=_D, bbox_inches="tight")
    plt.close(); print("  [gen] arch_baseline.png")

    # ════════════════════════════ ATTENTION v1 / v2 ═══════════════════════════
    fig, ax = plt.subplots(figsize=(14.8, 4.1), facecolor=_D)
    ax.set_facecolor(_D); ax.set_xlim(0,14.8); ax.set_ylim(0,4.1); ax.axis("off")

    ax.text(7.1, 3.93, "CNN + Multi-Head Self-Attention — SelvaSonic-ML    714,987 parámetros",
            ha="center", fontsize=13, color=_W, fontweight="bold")
    ax.text(7.1, 3.62,
            "v2 añade:  CrossEntropyLoss( weight=class_weights, label_smoothing=0.1 )"
            "  →  penaliza más errores en clases raras  (Chordeiles peso = 5.44 ×)",
            ha="center", fontsize=8.5, color=_L, style="italic")

    sbox(ax, 0.1, Y, 1.05, HB, _C,
         [("Input",9.5,_C,True),("[B,1,128,216]",7.5,_S,False),("Mel-spec",7,_M,False)])
    arr(ax, 1.15, 1.35, Y)

    _cb2 = [
        (1.35,2,_G,[("Conv 1",8.5,_G,True),("32 ch",7.5,_S,False),("BN·ReLU·Pool",6.5,_M,False)]),
        (2.62,2,_G,[("Conv 2",8.5,_G,True),("64 ch",7.5,_S,False),("BN·ReLU·Pool",6.5,_M,False)]),
        (3.89,2,_G,[("Conv 3",8.5,_G,True),("128 ch",7.5,_S,False),("BN·ReLU·Pool",6.5,_M,False)]),
        (5.16,2,_G,[("Conv 4",8.5,_G,True),("256 ch",7.5,_S,False),("BN·ReLU·Pool",6.5,_M,False)]),
    ]
    for i,(x,ns,col,lines) in enumerate(_cb2):
        fmap(ax, x, Y, 0.93, HF, ns, col, lines)
        xe = x+0.93+(ns-1)*0.068
        if i<3: arr(ax, xe+0.02, xe+0.14, Y)

    brace(ax, 1.35, 6.35, 0.42, "Feature Extraction (idéntico al Baseline)")

    xe4b = 5.16+0.93+0.068  # ≈ 6.158
    arr(ax, 6.18, 6.38, Y)
    sbox(ax, 6.38, Y, 1.0, 1.15, _M,
         [("Reshape",8.5,_M,True),("[B, 256, T]",7.5,_S,False),("tokens espaciales",7,_M,False)])
    arr(ax, 7.38, 7.58, Y)

    ax.add_patch(FancyBboxPatch((7.58,Y-HB/2-0.2),1.75,HB+0.4,
        boxstyle="round,pad=0.07",facecolor="#1C0910",edgecolor=_R,lw=3.0,zorder=3))
    _txt(ax, 7.58+0.875, Y,
         [("MHSA",12,_R,True),("4 heads · d_head = 64",8.5,_S,False),
          ("Layer Norm",7.5,_M,False),("Dropout  0.1",7,_M,False)])
    arr(ax, 9.33, 9.53, Y)

    brace(ax, 6.38, 9.53, 0.42, "Attention Module", _R)

    sbox(ax, 9.53, Y, 1.1, 1.15, _G,
         [("Mean Pool",9,_G,True),("[B, 256]",7.5,_S,False),("sobre tokens",7,_M,False)])
    arr(ax, 10.63, 10.83, Y)

    sbox(ax, 10.83, Y, 1.45, HB, _O,
         [("FC  256→128",9,_O,True),("ReLU",7.5,_S,False),("Dropout  0.3",7,_M,False)])
    arr(ax, 12.28, 12.48, Y)

    sbox(ax, 12.48, Y, 1.45, 1.65, _R,
         [("FC  128→11",9,_R,True),("LogSoftmax",7.5,_S,False),("11 clases",7,_M,False)])

    oh2 = HF/len(SPECIES)
    for k,sp in enumerate(SPECIES):
        yo = Y-HF/2+k*oh2
        fc2 = "#112F27" if sp!="no_ave" else "#0D2620"
        ec2 = _M if sp!="no_ave" else _C
        ax.add_patch(FancyBboxPatch((13.97,yo+0.01),0.8,oh2-0.02,
            boxstyle="square,pad=0",facecolor=fc2,edgecolor=ec2,lw=0.6,zorder=3))
        ax.text(14.37,yo+oh2/2,sp,ha="center",va="center",fontsize=5.5,color=_W,zorder=9)
    arr(ax, 13.93, 13.97, Y)

    brace(ax, 9.53, 14.77, 0.42, "Clasificación")

    plt.tight_layout(pad=0.18)
    plt.savefig(str(OUT_DIR/"arch_attention.png"), dpi=180, facecolor=_D, bbox_inches="tight")
    plt.close(); print("  [gen] arch_attention.png")


def gen_jungle_bg():
    """
    Compone fondo jungla: NEW JUNGLE_BG.PNG + overlay oscuro (#00100C).
    Blend: 40 % imagen + 60 % fondo → jungla visible con contraste suficiente para texto.
    """
    from PIL import Image as _PIL
    import numpy as np
    src = OUT_DIR / "NEW JUNGLE_BG.PNG"
    if not src.exists():
        print("  [skip jungle_bg] NEW JUNGLE_BG.PNG no encontrado en results/presentacion/"); return
    W, H = 1920, 1080
    try:
        resample = _PIL.Resampling.LANCZOS
    except AttributeError:
        resample = _PIL.LANCZOS
    img = _PIL.open(str(src)).convert("RGB").resize((W, H), resample)
    bg = (0, 16, 12)  # #00100C
    arr = np.array(img, dtype=np.float32)
    out = np.zeros((H, W, 3), dtype=np.float32)
    for c, v in enumerate(bg):
        out[:, :, c] = arr[:, :, c] * 0.40 + v * 0.60
    np.clip(out, 0, 255, out=out)
    _PIL.fromarray(out.astype(np.uint8), "RGB").save(str(JUNGLE_BG_COMPOSED))
    print(f"  [gen] jungle_bg.png ({W}×{H}) <- NEW JUNGLE_BG.PNG")


def gen_mel():
    try:
        import librosa, librosa.display
    except ImportError:
        print("  [skip mel] librosa no disponible"); return
    import matplotlib; matplotlib.use("Agg")
    import matplotlib.pyplot as plt, numpy as np
    d = ROOT / "data/raw/Crypturellus_cinereus"
    mp3s = sorted(d.glob("*.mp3"))
    if not mp3s: print("  [skip mel] sin mp3"); return
    mp3 = mp3s[0]
    y, sr = librosa.load(str(mp3), sr=22050, duration=5.0)
    S    = librosa.feature.melspectrogram(y=y, sr=sr, n_fft=2048, hop_length=512,
                                           n_mels=128, fmin=50, fmax=11025)
    Sdb  = librosa.power_to_db(S, ref=np.max)
    fig, ax = plt.subplots(figsize=(11, 4), facecolor="#00100C")
    ax.set_facecolor("#00100C")
    im = librosa.display.specshow(Sdb, sr=sr, hop_length=512, x_axis="time",
                                   y_axis="mel", fmin=50, fmax=11025, cmap="magma", ax=ax)
    ax.set_xlabel("Tiempo (s)", color="#EFFCF5", fontsize=10)
    ax.set_ylabel("Frecuencia (Hz)", color="#EFFCF5", fontsize=10)
    ax.tick_params(colors="#B5D3C3")
    ax.set_title(f"Crypturellus cinereus · 5 s · {mp3.stem}", color="#21D4D4", fontsize=11)
    for sp in ax.spines.values(): sp.set_color("#7EAA95"); sp.set_alpha(0.3)
    cb = plt.colorbar(im, ax=ax, format="%+2.0f dB")
    cb.ax.yaxis.set_tick_params(color="#B5D3C3")
    plt.setp(plt.getp(cb.ax.axes, "yticklabels"), color="#B5D3C3")
    plt.tight_layout()
    plt.savefig(IMGS["mel_ejemplo"], dpi=200, facecolor="#00100C", bbox_inches="tight")
    plt.close(); print("  [gen] mel_ejemplo.png")


# ── Helpers python-pptx ───────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SH_R  = 1   # RECTANGLE
SH_RR = 5   # ROUNDED_RECTANGLE

def rgb(h):  return RGBColor.from_string(h)
def I(v):    return Inches(v)
def P(v):    return Pt(v)

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def add_sl(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(sl, prs, c=BG):
    s = sl.shapes.add_shape(SH_R, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = rgb(c); s.line.fill.background()
    e = s._element; e.getparent().remove(e); sl.shapes._spTree.insert(2, e)

def set_bg_jungle(sl, prs):
    """Fondo jungla compuesto (jungle_bg.png) — fallback a sólido BG si falta la imagen."""
    if JUNGLE_BG_COMPOSED.exists():
        pic = sl.shapes.add_picture(str(JUNGLE_BG_COMPOSED), 0, 0,
                                    prs.slide_width, prs.slide_height)
        e = pic._element; e.getparent().remove(e); sl.shapes._spTree.insert(2, e)
    else:
        set_bg(sl, prs)

def tx(sl, text, x, y, w, h, sz, col, bold=False, al=PP_ALIGN.LEFT, it=False):
    s  = sl.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = s.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = al
    r  = p.add_run(); r.text = str(text)
    r.font.name = FONT; r.font.size = P(sz); r.font.bold = bold
    r.font.italic = it; r.font.color.rgb = rgb(col)
    return s

def tx_lines(sl, lines, x, y, w, h, sz=10, col=TS, bold=False):
    """lines: list of str | dict(text, sz, col, bold, sp_before)"""
    s  = sl.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = s.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if isinstance(item, str):
            r = p.add_run(); r.text = item
            r.font.name = FONT; r.font.size = P(sz)
            r.font.bold = bold; r.font.color.rgb = rgb(col)
        else:
            if item.get("sp"): p.space_before = P(item["sp"])
            r = p.add_run(); r.text = item.get("text", "")
            r.font.name = FONT
            r.font.size  = P(item.get("sz", sz))
            r.font.bold  = item.get("bold", bold)
            r.font.color.rgb = rgb(item.get("col", col))
    return s

def bx(sl, x, y, w, h, fill=PNL, border=None, bw=0.75, rnd=False):
    s = sl.shapes.add_shape(SH_RR if rnd else SH_R, I(x), I(y), I(w), I(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
    if border: s.line.color.rgb = rgb(border); s.line.width = P(bw)
    else: s.line.fill.background()
    return s

def hdr(sl, prs, title, sub=None, ty=0.18, tsz=22):
    tx(sl, title, 0.5, ty, 12.33, 0.55, tsz, TP, bold=True)
    if sub: tx(sl, sub, 0.5, ty + 0.52, 12.33, 0.34, 11.5, TM)

def ft(sl, prs, n):
    tx(sl, "SelvaSonic-ML · Aprendizaje Automático · UNAL Medellín",
       0.5, 7.15, 9.0, 0.3, 8, TM)
    tx(sl, f"{n:02d}", 12.3, 7.15, 0.53, 0.3, 8, TM, al=PP_ALIGN.RIGHT)

def add_img(sl, key, x, y, w, h=None):
    p = Path(IMGS.get(key, ""))
    if not p.exists(): print(f"  [miss] {key}"); return None
    kw = {"width": I(w), "height": I(h)} if h else {"width": I(w)}
    return sl.shapes.add_picture(str(p), I(x), I(y), **kw)

def add_framed_image(sl, key, x, y, w, h, frame_color=CYN, frame_pt=2.0):
    """
    Imagen envuelta en card blanco (#FAFAFA) con borde redondeado de color signature.
    La imagen se escala preservando aspect ratio (letterbox, centrada dentro del frame).
    """
    from PIL import Image as _PIL
    path = Path(IMGS.get(key, ""))
    if not path.exists():
        print(f"  [miss framed] {key}")
        bx(sl, x, y, w, h, fill=PNL, border=ORG, bw=1.0)
        tx(sl, f"TODO: {path.name}", x+0.1, y+h/2-0.15, w-0.2, 0.3, 9, ORG, al=PP_ALIGN.CENTER)
        return None, None

    inset = 0.06  # pulgadas de padding interior

    # Calcular dimensiones internas disponibles
    inner_w = w - 2 * inset
    inner_h = h - 2 * inset

    # Aspect ratio real de la imagen
    pil = _PIL.open(str(path))
    nat_w, nat_h = pil.size
    ratio = nat_h / nat_w

    # Escalar para encajar dentro del frame (letterbox, aspect ratio preservado)
    fit_w = inner_w
    fit_h = fit_w * ratio
    if fit_h > inner_h:
        fit_h = inner_h
        fit_w = fit_h / ratio

    # Centrar dentro del frame
    img_x = x + inset + (inner_w - fit_w) / 2
    img_y = y + inset + (inner_h - fit_h) / 2

    # 1. Card blanco con borde turquesa (queda DETRÁS de la imagen)
    card = sl.shapes.add_shape(SH_RR, I(x), I(y), I(w), I(h))
    card.fill.solid()
    card.fill.fore_color.rgb = rgb("FAFAFA")
    card.line.color.rgb = rgb(frame_color)
    card.line.width = P(frame_pt)
    try:
        card.adjustments[0] = min(0.05 / min(w, h), 0.2)
    except Exception:
        pass

    # 2. Imagen encima (python-pptx agrega formas en orden: la última queda arriba)
    pic = sl.shapes.add_picture(str(path), I(img_x), I(img_y), I(fit_w), I(fit_h))
    return card, pic

def set_cell_bg(cell, bg_h, fg_h, sz, bold=False, al=PP_ALIGN.CENTER):
    """Style table cell: background color + font."""
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for sf in tcPr.findall(qn("a:solidFill")): tcPr.remove(sf)
    sf  = etree.SubElement(tcPr, qn("a:solidFill"))
    clr = etree.SubElement(sf,  qn("a:srgbClr"))
    clr.set("val", bg_h)
    tf = cell.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = al
    r  = p.runs[0] if p.runs else p.add_run()
    r.font.name = FONT; r.font.size = P(sz)
    r.font.bold = bold; r.font.color.rgb = rgb(fg_h)

def make_table(sl, data, x, y, w, h, col_ws=None,
               hdr_bg=PNL, hdr_fg=CYN, hdr_sz=9,
               row_bg=BG,  row_bg2=PNL2, row_fg=TS, row_sz=8.5,
               win_col=None):
    """data[0] = header row.  win_col = column index to highlight winner."""
    rows, cols = len(data), len(data[0])
    tbl = sl.shapes.add_table(rows, cols, I(x), I(y), I(w), I(h)).table
    if col_ws:
        total = I(w)
        for ci, frac in enumerate(col_ws):
            tbl.columns[ci].width = int(total * frac)
    for ri, row in enumerate(data):
        for ci, txt in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(txt) if txt is not None else ""
            is_hdr  = ri == 0
            is_win  = (win_col is not None and ci == win_col and ri > 0)
            bg_c    = hdr_bg if is_hdr else (CYN if is_win else (row_bg2 if ri % 2 == 1 else row_bg))
            fg_c    = hdr_fg if is_hdr else (BG  if is_win else row_fg)
            sz_c    = hdr_sz if is_hdr else row_sz
            al_c    = PP_ALIGN.CENTER if ci == 0 or ri == 0 else PP_ALIGN.CENTER
            set_cell_bg(cell, bg_c, fg_c, sz_c, bold=is_hdr or is_win, al=al_c)
    return tbl


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def s01(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    bx(sl, 9.4, 0.0, 3.93, 7.5, fill=PNL)
    tx(sl, "Machine Learning · presentación técnica",
       0.6, 0.65, 8.5, 0.3, 11, CYN)
    tx(sl, "SelvaSonic-ML", 0.6, 1.05, 9.0, 1.45, 60, TP, bold=True)
    tx(sl, "Clasificación bioacústica de aves amazónicas",
       0.6, 2.5, 8.7, 0.4, 19, TS)
    tx(sl, "mediante CNN + Multi-Head Self-Attention",
       0.6, 2.9, 8.7, 0.4, 19, TS)
    tx(sl, ("Pipeline end-to-end · comparativa 3 modelos · evaluación "
            "multi-métrica (F1, AP, ECE, silhouette) · prueba en campo."),
       0.6, 3.5, 8.7, 0.45, 12, TM)
    tx(sl, "Laura Ruiz Arango · Jose Aldair Molina Méndez",
       0.6, 4.3, 8.7, 0.35, 12, TS)
    tx(sl, "Universidad Nacional de Colombia — Sede Medellín",
       0.6, 4.72, 8.7, 0.3, 10, TM)
    tx(sl, "Aprendizaje Automático · Prof. Alcides Montoya · Junio 2026",
       0.6, 5.04, 8.7, 0.3, 10, TM)
    return sl


def s02(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hd = ("1. Problema: escuchar la biodiversidad a escala",
          "Del monitoreo acústico pasivo a la clasificación automática de especies")
    hdr(sl, prs, *hd); ft(sl, prs, 2)
    cards = [("1,627", "especies registradas en Colombia"),
             ("3,319", "grabaciones en Amazonía colombiana"),
             ("10 + 1", "especies seleccionadas + no_ave")]
    for i, (val, lbl) in enumerate(cards):
        x0 = 0.5 + i * 4.2
        bx(sl, x0, 1.25, 4.0, 1.55, fill=PNL, border=CYN, bw=0.75)
        tx(sl, val, x0+0.1, 1.35, 3.8, 0.75, 42, CYN, bold=True, al=PP_ALIGN.CENTER)
        tx(sl, lbl, x0+0.1, 2.07, 3.8, 0.35, 9,  TM, al=PP_ALIGN.CENTER)
    bx(sl, 0.5, 3.0, 12.33, 1.12, fill=PNL2)
    tx(sl, ("La Amazonía colombiana concentra el 16 % de la diversidad mundial de aves. "
            "Monitorear su salud requiere identificación masiva de cantos, hoy dependiente de "
            "expertos humanos escasos. SelvaSonic-ML aborda este cuello de botella entrenando "
            "un clasificador bioacústico desde cero, especializado en fauna colombiana."),
       0.65, 3.07, 12.0, 0.98, 10.5, TS)
    bx(sl, 0.5, 4.25, 12.33, 1.4, fill=PNL, border=CYN)
    tx(sl, "Pregunta de investigación", 0.65, 4.3, 5.0, 0.32, 9, CYN, bold=True)
    tx(sl, ("¿Podemos aprender patrones espectro-temporales de vocalizaciones amazónicas "
            "con datos pequeños y desbalanceados, alcanzando rendimiento competitivo "
            "sin recurrir a modelos preentrenados?"),
       0.65, 4.65, 12.0, 0.9, 11.5, TP)
    return sl


def s03(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "2. Objetivos y antecedentes",
        "Diseño experimental con base científica clara")
    ft(sl, prs, 3)
    # Columna izquierda — objetivos
    bx(sl, 0.5, 1.08, 6.0, 5.7, fill=PNL)
    tx(sl, "Objetivos", 0.65, 1.15, 5.7, 0.35, 11, CYN, bold=True)
    tx_lines(sl, [
        {"text": "General:", "sz": 10, "col": TP, "bold": True},
        {"text": ("Diseñar, entrenar y evaluar un clasificador bioacústico end-to-end "
                  "basado en CNN + attention para 10 especies amazónicas + no_ave."),
         "sz": 9.5, "col": TS},
        {"text": " ", "sz": 4},
        {"text": "Específicos:", "sz": 10, "col": TP, "bold": True},
        {"text": "1.  Construir pipeline reproducible (Xeno-canto + ESC-50).",
         "sz": 9.5, "col": TS},
        {"text": "2.  Comparar arquitecturas: CNN baseline → +MHSA → +balanceo.",
         "sz": 9.5, "col": TS},
        {"text": "3.  Validar bajo cambio de distribución (audios externos).",
         "sz": 9.5, "col": TS},
        {"text": "4.  Documentar decisiones metodológicas con rigor científico.",
         "sz": 9.5, "col": TS},
        {"text": " ", "sz": 4},
        {"text": "Por qué sin preentrenamiento:", "sz": 10, "col": GLD, "bold": True},
        {"text": ("Modelos como YAMNet / AST se preentrenaron en AudioSet (etiquetas genéricas). "
                  "Aprender desde cero sobre fauna colombiana nos fuerza a construir "
                  "representaciones específicas — y a entender qué aprende el modelo realmente."),
         "sz": 9.5, "col": TS},
    ], 0.65, 1.52, 5.6, 5.1)
    # Columna derecha — tabla comparativa
    bx(sl, 6.65, 1.08, 6.18, 5.7, fill=PNL)
    tx(sl, "Antes vs. SelvaSonic-ML", 6.8, 1.15, 5.9, 0.35, 11, CYN, bold=True)
    tbl_data = [
        ["Aspecto",            "Proyecto anterior",   "SelvaSonic-ML"],
        ["Modelos",            "YAMNet, AST (pretrain.)", "CNN custom desde cero"],
        ["Dataset",            "3 audios Amazonas",   "~14 K clips train"],
        ["Clases",             "AudioSet genérico",   "10 sp. amazónicas"],
        ["Métricas",           "Cualitativas",        "F1, AP, ECE, silhouette"],
        ["Pipeline",           "Sin estructurar",     "Modular, testeable"],
        ["Interpretabilidad",  "Ninguna",             "Grad-CAM + t-SNE"],
    ]
    make_table(sl, tbl_data, 6.78, 1.55, 5.9, 4.95,
               col_ws=[0.3, 0.35, 0.35],
               hdr_bg=BG2, hdr_fg=CYN, hdr_sz=8.5,
               row_bg=PNL, row_bg2=BG2, row_fg=TS, row_sz=8,
               win_col=2)
    return sl


def s04(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "3. Estructura del repositorio",
        "Código modular, testeable, con type hints y config centralizada")
    ft(sl, prs, 4)
    # Diagrama de flujo: 5 nodos horizontales
    flow = [
        (CYN,  "data/",     "Xeno-canto + ESC-50\nraw → processed"),
        (GRN,  "src/",      "audio_io · transforms\ndataset · model · loss\nclass_weights · inference"),
        (GLD,  "scripts/",  "calcular_class_weights\ntrain · evaluate"),
        (ORG,  "notebooks/","EDA → baseline → attention\n→ comparativa → Grad-CAM\n→ embeddings → demo"),
        (CYN,  "results/",  "runs/ · comparativa/\ngradcam/ · embeddings/\ntest_externos/"),
    ]
    for i, (col, title, desc) in enumerate(flow):
        x0 = 0.5 + i * 2.56
        bx(sl, x0, 1.1, 2.3, 2.1, fill=PNL, border=col, bw=0.75)
        tx(sl, title, x0+0.1, 1.18, 2.1, 0.35, 11, col, bold=True)
        tx(sl, desc,  x0+0.1, 1.56, 2.1, 1.45, 8.5, TS)
        if i < 4:
            tx(sl, "→", x0+2.3, 1.9, 0.26, 0.35, 16, TM, al=PP_ALIGN.CENTER)
    # Principios de diseño
    bx(sl, 0.5, 3.35, 12.33, 1.15, fill=PNL2)
    tx(sl, "Decisiones de diseño:", 0.65, 3.42, 4.0, 0.3, 10, CYN, bold=True)
    tx_lines(sl, [
        "·  Config centralizada (src/config.py) — único lugar para todos los hiperparámetros.",
        "·  Semilla 42 fija en todos los runs — splits idénticos para comparación justa.",
        "·  Early stopping con patience configurable — detecta convergencia automáticamente.",
        "·  Logger TensorBoard por epoch — curvas reproducibles fuera de Colab.",
    ], 0.65, 3.74, 12.0, 1.6, 9.5, TS)
    # Módulos clave
    bx(sl, 0.5, 4.65, 12.33, 1.35, fill=PNL, border=TM, bw=0.4)
    tx(sl, "Módulos clave — src/", 0.65, 4.7, 4.0, 0.3, 10, GLD, bold=True)
    tx(sl, ("model.py (SelvaSonicCNN / CNNAttention)  ·  loss.py (CrossEntropy + class_weights + label_smoothing)  ·  "
            "class_weights.py (cómputo de pesos inversos)  ·  transforms.py (Mel-espectrogramas)  ·  "
            "inference.py (pipeline de inferencia con umbral de confianza)  ·  logger.py (TensorBoard)"),
       0.65, 5.02, 12.0, 0.9, 9.5, TS)
    tx(sl, "Fuente: github.com/LauraRuizA09/SelvaSonic-ML",
       0.5, 6.7, 12.33, 0.28, 8, TM)
    return sl


def s05(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "4. Dataset: 10 especies amazónicas + clase no_ave",
        "Diseño con diversidad taxonómica (8 familias) · reto intra-género (2 Crypturellus)")
    ft(sl, prs, 5)
    headers = ["#", "Especie", "Familia", "Grabaciones\n(XC)", "Clips train",  "Weight"]
    rows = [[str(i+1)] + [sp, fam, str(grab), f"{clips:,}", wt]
            for i, (sp, fam, grab, clips, wt) in enumerate(SPECIES)]
    rows.append(["—", "no_ave (ESC-50, sin animales)", "—", "~1,600 archivos", "2,240", "0.5833"])
    data = [headers] + rows
    make_table(sl, data, 0.5, 1.08, 12.33, 5.45,
               col_ws=[0.045, 0.27, 0.16, 0.12, 0.12, 0.085],
               hdr_bg=BG2, hdr_fg=CYN, hdr_sz=9,
               row_bg=BG, row_bg2=PNL2, row_fg=TS, row_sz=8.5)
    tx(sl, ("Criterios de selección: ≥14 grabaciones · diversidad taxonómica (8 familias) · "
            "reto intra-género (Crypturellus) · exclusión de 'Mystery mystery' · "
            "no_ave de ESC-50 excluyendo categoría Animals."),
       0.5, 6.64, 12.33, 0.4, 8, TM)
    return sl


def s06(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "5. Desbalance severo de clases",
        "no_ave no es la más abundante — Trogon viridis lidera · las minoritarias motivan class weights")
    ft(sl, prs, 6)
    add_img(sl, "dist_clases", 0.5, 1.08, 8.8)
    # Callout derecho
    bx(sl, 9.45, 1.08, 3.38, 5.6, fill=PNL, border=CYN, bw=0.6)
    tx_lines(sl, [
        {"text": "Estadísticas clave", "sz": 10, "col": CYN, "bold": True},
        {"text": " ", "sz": 3},
        {"text": "Total train:", "sz": 9, "col": GLD, "bold": True},
        {"text": "14,372 clips", "sz": 11, "col": TP, "bold": True},
        {"text": " ", "sz": 3},
        {"text": "Clase más frecuente:", "sz": 9, "col": TS},
        {"text": "Trogon viridis · 2,784", "sz": 10, "col": TP, "bold": True},
        {"text": " ", "sz": 3},
        {"text": "Clase más rara:", "sz": 9, "col": TS},
        {"text": "Chordeiles · 240", "sz": 10, "col": RED, "bold": True},
        {"text": " ", "sz": 3},
        {"text": "Ratio max/min:", "sz": 9, "col": TS},
        {"text": "11.6×  (Trogon / Chordeiles)", "sz": 10, "col": ORG, "bold": True},
        {"text": " ", "sz": 3},
        {"text": "Peso máximo:", "sz": 9, "col": TS},
        {"text": "5.44 (Chordeiles)", "sz": 10, "col": CYN, "bold": True},
        {"text": " ", "sz": 3},
        {"text": "Decisión:", "sz": 9, "col": GLD, "bold": True},
        {"text": "atacar con class weights en la loss antes de tunear arquitectura.",
         "sz": 9, "col": TS},
    ], 9.6, 1.18, 3.1, 5.4)
    return sl


def s07(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "6. Pipeline de preprocesamiento",
        "Carga → Segmentación → Mel-espectrograma → Normalización → Tensor")
    ft(sl, prs, 7)
    steps = [
        ("Carga", "librosa\nsr=22 050 Hz\nmono"),
        ("Segmentación", "clips 5 s\nsolape 50 %\n(solo train)"),
        ("Mel FFT", "n_fft=2 048\nhop=512\nn_mels=128"),
        ("Filtros", "fmin=50 Hz\nfmax=11 025 Hz"),
        ("Tensor", "dB norm.\nshape: [1,128,~216]"),
    ]
    for i, (title, desc) in enumerate(steps):
        x0 = 0.45 + i * 2.57
        bx(sl, x0, 1.1, 2.3, 2.25, fill=PNL, border=CYN, bw=0.75, rnd=True)
        tx(sl, title, x0+0.1, 1.2,  2.1, 0.35, 11, CYN, bold=True, al=PP_ALIGN.CENTER)
        tx(sl, desc,  x0+0.1, 1.58, 2.1, 1.55, 8.5, TS, al=PP_ALIGN.CENTER)
        if i < 4:
            tx(sl, "→", x0+2.3, 2.0, 0.27, 0.4, 16, TM, al=PP_ALIGN.CENTER)
    # Tabla de hiperparámetros
    bx(sl, 0.5, 3.55, 12.33, 0.35, fill=BG2)
    tx(sl, "Hiperparámetros del pipeline (src/config.py — fuente única de verdad)",
       0.65, 3.6, 12.0, 0.28, 9, GLD, bold=True)
    tbl = [
        ["Parámetro",    "Valor",  "Justificación"],
        ["SAMPLE_RATE",  "22 050 Hz", "Nyquist cubre hasta 11 025 Hz > rango aves (500 Hz–10 kHz)"],
        ["CLIP_DURATION","5.0 s",  "Estándar en BirdCLEF/BirdNET; captura cantos completos"],
        ["N_FFT",        "2 048",  "Δf≈10.8 Hz/bin — suficiente para distinguir tonos de aves"],
        ["HOP_LENGTH",   "512",    "Δt≈23 ms/frame, overlap 75 % — captura trinos rápidos"],
        ["N_MELS",       "128",    "Estándar YAMNet/PANNs; equilibra detalle vs. cómputo"],
        ["CLIP_OVERLAP",  "0.5 (solo train)", "Data augmentation temporal honesto — no en val/test"],
    ]
    make_table(sl, tbl, 0.5, 3.93, 12.33, 3.0,
               col_ws=[0.18, 0.14, 0.68],
               hdr_bg=PNL, hdr_fg=CYN, hdr_sz=8.5,
               row_bg=BG, row_bg2=PNL2, row_fg=TS, row_sz=8.5,
               win_col=None)
    return sl


def s08(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "6b. Del audio crudo al Mel-espectrograma",
        "Cada cuadro tiempo-frecuencia recibe color según energía en dB")
    ft(sl, prs, 8)
    added = add_img(sl, "mel_ejemplo", 0.5, 1.08, 12.33)
    if not added:
        bx(sl, 0.5, 1.08, 12.33, 4.5, fill=PNL, border=TM, bw=0.4)
        tx(sl, "TODO: regenerar mel_ejemplo.png con venv Python (librosa disponible en venv)",
           1.0, 2.8, 10.0, 0.5, 11, ORG, al=PP_ALIGN.CENTER)
    bx(sl, 0.5, 5.8, 12.33, 1.05, fill=PNL2)
    tx_lines(sl, [
        {"text": "Lectura del espectrograma:", "sz": 10, "col": CYN, "bold": True},
        {"text": ("Eje X = tiempo (s)  ·  Eje Y = frecuencia (escala Mel)  ·  "
                  "Color (magma) = energía en dB relativo al pico (rojo = alta energía).  "
                  "Las CNN aprenden filtros que detectan patrones espectrales como el silbo "
                  "descendente de Crypturellus o el graznido pulsado de Ramphastos."),
         "sz": 9.5, "col": TS},
    ], 0.65, 5.85, 12.0, 0.95)
    return sl


def s09(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "7. Augmentation y control experimental",
        "Solo en TRAIN — nunca en val/test para evitar data leakage")
    ft(sl, prs, 9)
    augs = [
        (CYN,  "Time Stretch",   "rate ∈ [0.8, 1.25]",
         "P=0.50", "Simula aves cantando a velocidades distintas (excitación, fatiga). "
                   "±25 % sin artefactos del phase vocoder."),
        (GLD,  "Pitch Shift",    "±2 semitonos",
         "P=0.50", "Cubre variación natural entre individuos (machos, hembras, juveniles). "
                   "f' = f × 2^(n/12); ±2 semitonos → [0.89×, 1.12×]."),
        (GRN,  "Add Noise",      "SNR ∈ [10, 30] dB",
         "P=0.70", "Simula condiciones de campo. SNR=30 dB silencioso, "
                   "SNR=10 dB lluvia/ambiente. Más alto P porque es la condición más realista."),
    ]
    for i, (col, title, param, prob, desc) in enumerate(augs):
        x0 = 0.5 + i * 4.2
        bx(sl, x0, 1.1, 4.0, 3.55, fill=PNL, border=col, bw=0.75, rnd=True)
        tx(sl, title, x0+0.15, 1.18, 3.7, 0.35, 12, col, bold=True)
        tx(sl, param, x0+0.15, 1.55, 3.7, 0.3,  10, TP)
        tx(sl, prob,  x0+0.15, 1.85, 3.7, 0.28,  9, TM)
        tx(sl, desc,  x0+0.15, 2.2,  3.7, 2.2,   9, TS)
    bx(sl, 0.5, 4.82, 12.33, 1.1, fill=PNL2)
    tx_lines(sl, [
        {"text": "Principios del control experimental:", "sz": 10, "col": GLD, "bold": True},
        {"text": ("·  Augmentation SOLO en train splits — val y test usan audio crudo sin modificar.  "
                  "·  La augmentation es temporal: cada epoch el mini-batch puede ver versiones distintas del mismo audio.  "
                  "·  Short Audio Strategy = 'wrap' (loop circular) en lugar de zero-padding "
                  "— evita que el modelo aprenda 'borde de silencio' como feature artefacto."),
         "sz": 9.5, "col": TS},
    ], 0.65, 4.9, 12.0, 0.95)
    return sl


def s10(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "8a. Arquitectura 1: Baseline CNN",
        "Punto de partida de la progresión metodológica · Input → 4 Conv Blocks → GAP → FC → 11 clases")
    ft(sl, prs, 10)
    # Diagrama de arquitectura generado con matplotlib
    add_img(sl, "arch_baseline", 0.5, 0.88, 12.33)
    # Stats bar
    bx(sl, 0.5, 4.32, 12.33, 1.05, fill=BG2)
    stats = [("422,635", "parámetros totales", CYN),
             ("27 / 30",  "epochs (early stop.)", GRN),
             ("0.6322",   "test accuracy",        GLD),
             ("0.509",    "macro F1",              ORG)]
    for j, (val, lbl, col) in enumerate(stats):
        x0 = 0.65 + j * 3.07
        tx(sl, val, x0, 4.40, 3.0, 0.50, 30, col, bold=True)
        tx(sl, lbl, x0, 4.95, 3.0, 0.30, 10, TM)
    # Decisiones de diseño
    bx(sl, 0.5, 5.48, 12.33, 1.72, fill=PNL)
    tx_lines(sl, [
        {"text": "Decisiones de diseño:", "sz": 11, "col": CYN, "bold": True},
        {"text": ("·  Progresión 32→64→128→256 (estándar VGGNet): cada bloque extrae representaciones más abstractas."),
         "sz": 10.5, "col": TS},
        {"text": ("·  Global Average Pooling en lugar de Flatten: elimina dependencia de la dimensión temporal de entrada."),
         "sz": 10.5, "col": TS},
        {"text": ("·  Dropout 0.3 conservador: con datasets pequeños, valores mayores degradan el flujo de gradientes."),
         "sz": 10.5, "col": TS},
    ], 0.68, 5.56, 12.0, 1.58)
    return sl


def s11(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "8b. Arquitectura 2 y 3: + Multi-Head Self-Attention (+v2 class weights)",
        "Un solo cambio estructural sobre Baseline: insertar MHSA entre feature extractor y clasificador")
    ft(sl, prs, 11)
    # Diagrama de arquitectura generado con matplotlib
    add_img(sl, "arch_attention", 0.5, 0.88, 12.33)
    # Stats bar — v1 vs v2
    bx(sl, 0.5, 4.32, 12.33, 1.05, fill=BG2)
    stats = [("714,987",  "parámetros (+69 % vs baseline)", CYN),
             ("17 / 30",  "v1: epochs (converge antes)",    GRN),
             ("0.7036",   "v2: test accuracy (+7.1 pp)",    GLD),
             ("0.625",    "v2: macro F1 (+11.6 pp)",        ORG)]
    for j, (val, lbl, col) in enumerate(stats):
        x0 = 0.65 + j * 3.07
        tx(sl, val, x0, 4.40, 3.0, 0.50, 30, col, bold=True)
        tx(sl, lbl, x0, 4.95, 3.0, 0.30, 10, TM)
    # Por qué MHSA mejora
    bx(sl, 0.5, 5.48, 12.33, 1.72, fill=PNL)
    tx_lines(sl, [
        {"text": "Por qué MHSA mejora sobre el Baseline:", "sz": 11, "col": CYN, "bold": True},
        {"text": ("·  Cada frame temporal 'consulta' todos los demás frames y pondera cuáles son más relevantes: "
                  "captura la estructura temporal de trinos y series de notas."),
         "sz": 10.5, "col": TS},
        {"text": ("·  Converge en 17 épocas vs 27 del baseline: la atención encuentra patrones discriminativos más eficientemente."),
         "sz": 10.5, "col": TS},
        {"text": ("·  v2 añade class_weights a la loss — mismo número de parámetros, solo cambia la función de costo."),
         "sz": 10.5, "col": GLD},
    ], 0.68, 5.56, 12.0, 1.58)
    return sl


def s12(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "9. Modelo final v2: atacando el desbalance con class weights",
        "Misma arquitectura que v1 · mismo label smoothing α=0.1 · cambia SOLO una cosa: los pesos")
    ft(sl, prs, 12)
    # Fórmula
    bx(sl, 0.5, 1.08, 7.0, 2.85, fill=PNL2)
    tx(sl, "Fórmula de class weights:", 0.65, 1.15, 6.7, 0.32, 10, CYN, bold=True)
    bx(sl, 0.65, 1.5, 6.7, 0.8, fill=BG, border=CYN, bw=0.4)
    tx(sl, "w_c  =  N_total  /  (K  ×  N_c)",
       0.75, 1.6, 6.5, 0.55, 18, TP, bold=True, al=PP_ALIGN.CENTER)
    tx_lines(sl, [
        {"text": "Donde:", "sz": 9, "col": TM, "bold": True},
        {"text": "N_total = 14,372  (total clips de entrenamiento)", "sz": 9.5, "col": TS},
        {"text": "K = 11  (número de clases)", "sz": 9.5, "col": TS},
        {"text": "N_c = clips de la clase c en train", "sz": 9.5, "col": TS},
        {"text": " ", "sz": 3},
        {"text": "Aplicado como: CrossEntropyLoss(weight=w, label_smoothing=0.1)",
         "sz": 9.5, "col": GLD},
        {"text": " ", "sz": 3},
        {"text": "Efecto: un error en Chordeiles (peso 5.44) penaliza 11.6× más",
         "sz": 9.5, "col": TS},
        {"text": "que un error en Trogon viridis (peso 0.47).",
         "sz": 9.5, "col": TS},
    ], 0.65, 2.33, 6.7, 2.5)
    # Tabla de pesos reales
    bx(sl, 7.65, 1.08, 5.18, 2.85, fill=PNL)
    tx(sl, "Pesos reales (data/class_weights.pt)", 7.8, 1.15, 4.9, 0.32, 9.5, CYN, bold=True)
    w_data = [
        ["Clase",              "N_train", "Peso"],
        ["Trogon viridis",     "2,784",   "0.4693 ⬇ mín"],
        ["Crypturellus cin.",  "2,656",   "0.4919"],
        ["no_ave",             "2,240",   "0.5833"],
        ["Rupornis mag.",       "302",    "4.3263"],
        ["Chordeiles pus.",     "240",    "5.4439 ⬆ máx"],
    ]
    make_table(sl, w_data, 7.73, 1.5, 4.95, 2.35,
               col_ws=[0.45, 0.3, 0.25],
               hdr_bg=BG2, hdr_fg=CYN, hdr_sz=8,
               row_bg=PNL, row_bg2=BG2, row_fg=TS, row_sz=8.5,
               win_col=2)
    # Nota sobre label smoothing
    bx(sl, 0.5, 4.1, 12.33, 0.9, fill=PNL2)
    tx(sl, "Nota sobre label smoothing (presente en los 3 modelos):", 0.65, 4.16, 7.0, 0.3, 9, GLD, bold=True)
    tx(sl, ("Los 3 modelos usan label_smoothing=0.1. La decisión fue aplicarlo desde baseline "
            "para tener una comparación justa. La hipótesis original era reservarlo para v2, "
            "pero mantenerlo en todos aísla el efecto de class weights con mayor rigor."),
       0.65, 4.47, 12.0, 0.48, 9.5, TS)
    # Decisión metodológica
    bx(sl, 0.5, 5.15, 12.33, 1.5, fill=PNL, border=CYN, bw=0.75)
    tx(sl, "Decisión metodológica clave", 0.65, 5.22, 5.0, 0.32, 10, CYN, bold=True)
    tx(sl, ("NO se cambia la arquitectura entre v1 y v2. "
            "El único cambio es use_class_weights=True en la loss. "
            "Cualquier mejora en métricas se atribuye exclusivamente al balanceo, "
            "no a capacidad adicional del modelo. Esto es un experimento controlado."),
       0.65, 5.57, 12.0, 1.0, 10.5, TP)
    return sl


def s13(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "10. Protocolo de entrenamiento",
        "Mismas semillas · mismo split · mismo hardware → comparación justa entre los 3 modelos")
    ft(sl, prs, 13)
    hp_data = [
        ["Hiperparámetro",    "Valor",            "Justificación"],
        ["Optimizer",         "AdamW",             "Estándar moderno; regularización implícita vía weight decay"],
        ["Learning rate",     "0.001",             "Valor de arranque; CosineAnnealingLR decae a 0"],
        ["LR Scheduler",      "CosineAnnealingLR", "Decaimiento suave; evita oscilaciones de SGD con momentum"],
        ["Weight decay",      "1e-4",              "Regularización L2 implícita en AdamW"],
        ["Label smoothing",   "α = 0.1",           "En los 3 modelos; reduce overconfidence (Szegedy et al. 2016)"],
        ["Batch size",        "32",                "Trade-off memoria GPU T4 / ruido en gradientes"],
        ["Max epochs",        "30",                "Límite superior; early stopping actúa antes en todos los runs"],
        ["Early stopping",    "patience=5, Δ=0.001","Detiene cuando val_acc no mejora ≥0.001 en 5 épocas"],
        ["Semilla",           "42",                "Fija splits, init de pesos y augmentation para reproducibilidad"],
        ["Hardware",          "NVIDIA T4 (Colab)", "GPU gratuita de Google Colab; reproducible con cuenta estándar"],
    ]
    make_table(sl, hp_data, 0.5, 1.08, 12.33, 4.65,
               col_ws=[0.22, 0.2, 0.58],
               hdr_bg=BG2, hdr_fg=CYN, hdr_sz=8.5,
               row_bg=BG, row_bg2=PNL2, row_fg=TS, row_sz=8.5)
    # Training times
    bx(sl, 0.5, 5.9, 12.33, 0.95, fill=PNL)
    tx(sl, "Tiempos de entrenamiento reales (GPU T4 · medidos en history.json):",
       0.65, 5.96, 8.0, 0.3, 9.5, CYN, bold=True)
    times = [("Baseline", "~286 min (27 ep.)"), ("Attention v1", "~181 min (17 ep.)"), ("Attention v2", "~270 min (25 ep.)")]
    for i, (name, t_str) in enumerate(times):
        x0 = 0.65 + i * 4.1
        tx(sl, name,  x0, 6.3,  3.9, 0.28, 9, TM)
        tx(sl, t_str, x0, 6.57, 3.9, 0.28, 11, TP, bold=True)
    return sl


def s14(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "11. Resultados: progresión metodológica",
        "Cada cambio aporta una mejora medible en test (mismo split, misma semilla)")
    ft(sl, prs, 14)
    tbl = [
        ["Métrica",            "Baseline",    "Attention v1",  "Attention v2 ⭐"],
        ["val_acc",            "0.7419",      "0.7829",        "0.7653"],
        ["test_acc",           "0.6322",      "0.6873",        "0.7036"],
        ["Macro F1",           "0.509",       "0.599",         "0.625"],
        ["Macro Precision",    "0.525",       "0.606",         "0.638"],
        ["Macro Recall",       "0.536",       "0.618",         "0.662"],
        ["AP medio cl. raras", "0.456",       "0.573",         "0.664"],
        ["Gap val-test",       "0.1097",      "0.0956",        "0.0617"],
        ["ECE (calibración)",  "0.050",       "0.039",         "0.145 ⚠"],
        ["Parámetros",         "422,635",     "714,987",       "714,987"],
        ["Epochs (/ 30)",      "27",          "17",            "25"],
    ]
    make_table(sl, tbl, 0.5, 1.08, 12.33, 5.0,
               col_ws=[0.26, 0.22, 0.22, 0.30],
               hdr_bg=BG2, hdr_fg=CYN, hdr_sz=9,
               row_bg=BG, row_bg2=PNL2, row_fg=TS, row_sz=9,
               win_col=3)
    bx(sl, 0.5, 6.25, 12.33, 0.85, fill=PNL, border=CYN, bw=0.5)
    tx(sl, ("Lectura clave: v2 no solo mejora test_acc (+7.1 pp vs baseline). "
            "El gap val-test bajó de 0.110 → 0.062 (−44 %, mejor generalización). "
            "AP en clases raras mejoró +21 pp. ECE sube en v2 (trade-off conocido con class weights — ver slide 17). "
            "AP medio clases raras = AP de Frederickena, Rupornis, Chordeiles, Glaucidium, Celeus."),
       0.65, 6.32, 12.0, 0.72, 9, TS)
    return sl


def s15(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "11b. Curvas de entrenamiento — Attention v2",
        "Convergencia progresiva · early stopping en epoch 25 de 30")
    ft(sl, prs, 15)
    # curvas_v2: 2684x731 px (ratio 0.272 — imagen muy panorámica)
    add_framed_image(sl, "curvas_v2", 0.5, 1.08, 8.7, 2.6)
    bx(sl, 9.35, 1.08, 3.48, 5.6, fill=PNL)
    tx_lines(sl, [
        {"text": "Lecturas clave", "sz": 10, "col": CYN, "bold": True},
        {"text": " ", "sz": 4},
        {"text": "1.  Loss train/val descienden en paralelo → no overfitting severo.",
         "sz": 9.5, "col": TS},
        {"text": " ", "sz": 3},
        {"text": "2.  val_acc mejor en epoch 15 → early stopping se activa en ep. 25 "
                 "(patience=5, no mejoró 0.001).",
         "sz": 9.5, "col": TS},
        {"text": " ", "sz": 3},
        {"text": "3.  La curva de pérdida es más ruidosa que v1: class weights amplifican "
                 "el gradiente de clases raras → más varianza por batch.",
         "sz": 9.5, "col": TS},
        {"text": " ", "sz": 3},
        {"text": "4.  Gap train/val en loss (v2 usa class weights en train pero no en val) "
                 "→ val_loss parece mayor, pero test_acc es el más alto.",
         "sz": 9.5, "col": GLD},
    ], 9.5, 1.2, 3.2, 5.4)
    return sl


def s16(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "12. Análisis de errores: matrices de confusión",
        "v2 fortalece la diagonal · especialmente en clases raras")
    ft(sl, prs, 16)
    # confusion_3: 2949x916 px (ratio 0.311 — panorámica)
    add_framed_image(sl, "confusion_3", 0.5, 1.08, 12.33, 4.0)
    bx(sl, 0.5, 5.22, 12.33, 0.9, fill=PNL2)
    tx_lines(sl, [
        {"text": "·  Diagonal más fuerte en v2 — especialmente Crypturellus undulatus (F1: 0.264 → 0.526) "
                 "y Rupornis magnirostris (F1: 0.053 → 0.250).",
         "sz": 9.5, "col": TS},
        {"text": "·  Confusiones residuales son taxonómicamente razonables: "
                 "Crypturellus cinereus ↔ C. undulatus (mismo género, canto similar).",
         "sz": 9.5, "col": TS},
    ], 0.65, 5.28, 12.0, 0.78)
    return sl


def s17(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "13. Trade-off honesto: class weights mejoran F1, cuestan ECE",
        "Los rebalanceos de clase tienen un costo conocido en calibración (Cao et al. 2019)")
    ft(sl, prs, 17)
    # Valores ECE — 3 tarjetas
    ece_cards = [
        ("Baseline",      "0.050", TM,   "calibración media"),
        ("Attention v1",  "0.039", GRN,  "mejor calibración ⭐"),
        ("Attention v2",  "0.145", ORG,  "mayor ECE ⚠ (class weights)"),
    ]
    for i, (name, val, col, note) in enumerate(ece_cards):
        x0 = 0.5 + i * 4.2
        bx(sl, x0, 1.1, 4.0, 1.65, fill=PNL, border=col, bw=0.75)
        tx(sl, name,  x0+0.1, 1.17, 3.8, 0.32, 10, col, bold=True, al=PP_ALIGN.CENTER)
        tx(sl, f"ECE = {val}", x0+0.1, 1.48, 3.8, 0.5, 26, col, bold=True, al=PP_ALIGN.CENTER)
        tx(sl, note,  x0+0.1, 2.0,  3.8, 0.3,   9, TM, al=PP_ALIGN.CENTER)
    # calibracion: 2687x763 px (ratio 0.284 — panorámica)
    add_framed_image(sl, "calibracion", 0.5, 2.9, 7.5, 2.3)
    # Interpretación
    bx(sl, 8.15, 2.9, 4.68, 3.8, fill=PNL2)
    tx_lines(sl, [
        {"text": "¿Qué es ECE?", "sz": 10, "col": CYN, "bold": True},
        {"text": ("Expected Calibration Error mide si la confianza del modelo "
                  "coincide con su accuracy real. ECE=0 es perfecto."),
         "sz": 9, "col": TS},
        {"text": " ", "sz": 4},
        {"text": "¿Por qué v2 tiene ECE mayor?", "sz": 10, "col": ORG, "bold": True},
        {"text": ("Los class weights fuerzan al modelo a ser más confiado en clases "
                  "minoritarias para compensar su escasez. Esto mejora F1 macro "
                  "(cada clase pesa igual), pero deteriora ECE global "
                  "(dominado por clases mayoritarias)."),
         "sz": 9, "col": TS},
        {"text": " ", "sz": 4},
        {"text": "Decisión metodológica:", "sz": 10, "col": GLD, "bold": True},
        {"text": ("Priorizar F1 macro porque en monitoreo bioacústico interesa más "
                  "detectar especies raras que tener confianza calibrada en las dominantes."),
         "sz": 9, "col": TS},
        {"text": " ", "sz": 4},
        {"text": "Mitigación futura: temperature scaling post-hoc sobre v2.", "sz": 9, "col": TM},
    ], 8.3, 3.0, 4.4, 3.65)
    return sl


def s18(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "14. Espacio de representaciones aprendido",
        "Silhouette score 20× mejor: las clases se separan en el espacio de embeddings 256-d")
    ft(sl, prs, 18)
    # Stats callout — 3 tarjetas
    sil_cards = [
        ("Baseline",     "+0.0084", RED,  "clases mezcladas"),
        ("Attention v1", "+0.1580", GRN,  "×19 mejora"),
        ("Attention v2", "+0.1686", CYN,  "×20 mejora ⭐"),
    ]
    for i, (name, val, col, note) in enumerate(sil_cards):
        x0 = 0.5 + i * 4.2
        bx(sl, x0, 1.1, 4.0, 1.55, fill=PNL, border=col, bw=0.6)
        tx(sl, name,  x0+0.1, 1.17, 3.8, 0.3, 9, col, bold=True, al=PP_ALIGN.CENTER)
        tx(sl, f"Silhouette = {val}", x0+0.1, 1.47, 3.8, 0.45, 22, col, bold=True, al=PP_ALIGN.CENTER)
        tx(sl, note,  x0+0.1, 1.94, 3.8, 0.28, 8.5, TM, al=PP_ALIGN.CENTER)
    # embeddings: 3392x2130 px (ratio 0.628 — letterbox)
    add_framed_image(sl, "embeddings", 0.5, 2.8, 12.33, 3.65)
    tx(sl, ("Nota: el silhouette score se calcula sobre el espacio original de 256 dimensiones (no sobre t-SNE). "
            "La visualización usa t-SNE (perplexity=30). Trustworthiness t-SNE ≈ 0.996 en los 3 modelos → "
            "la proyección 2D preserva bien la estructura del espacio real."),
       0.5, 6.58, 12.33, 0.42, 8, TM)
    return sl


def s19(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "15. Interpretabilidad: Grad-CAM sobre Attention v2",
        "El modelo focaliza en regiones tiempo-frecuencia con vocalización, no en silencios")
    ft(sl, prs, 19)
    # gradcam_cine: 1787x1099 px · gradcam_q: 1669x588 px (frame_pt=1.5)
    add_framed_image(sl, "gradcam_cine", 0.5, 1.08, 6.5,  4.2, frame_pt=1.5)
    add_framed_image(sl, "gradcam_q",   7.1, 1.08, 5.73, 4.2, frame_pt=1.5)
    bx(sl, 0.5, 5.85, 12.33, 1.1, fill=PNL2)
    tx_lines(sl, [
        {"text": "Grad-CAM proyecta los gradientes del último mapa de activación sobre el espectrograma de entrada. "
                 "Las regiones cálidas son las que más contribuyeron a la predicción correcta.", "sz": 9.5, "col": TS},
        {"text": ("Hallazgo cuantitativo: entropía media de los mapas CAM bajó de 0.952 (baseline) → 0.841 (v2). "
                  "Entropía más baja = atención más focalizada. "
                  "Centroide de frecuencia subió de 65.9 → 76.7 bins Mel: v2 mira más arriba en el espectro."),
         "sz": 9.5, "col": TS, "sp": 3},
    ], 0.65, 5.92, 12.0, 0.98)
    return sl


def s20(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "16. Prueba en campo: 6 audios externos del Amazonas (Puerto Nariño)",
        "Sin ground truth — caracteriza comportamiento bajo cambio de distribución")
    ft(sl, prs, 20)
    # ext_dist: 2991x3644 px (ratio 1.218 — letterbox vertical)
    add_framed_image(sl, "ext_dist", 0.5, 1.08, 8.5, 5.6)
    bx(sl, 9.1, 1.08, 3.73, 5.6, fill=PNL)
    tx_lines(sl, [
        {"text": "Métricas de campo", "sz": 10, "col": CYN, "bold": True},
        {"text": " ", "sz": 3},
        {"text": "Confianza media", "sz": 9, "col": TM},
        {"text": "Baseline:  0.40", "sz": 10, "col": TS},
        {"text": "Att. v1:   0.59  (más alta)", "sz": 10, "col": ORG},
        {"text": "Att. v2:   0.42  (más baja ⭐)", "sz": 10, "col": CYN},
        {"text": " ", "sz": 3},
        {"text": "% no_identificado (θ=0.6)", "sz": 9, "col": TM},
        {"text": "Baseline:  90.2 %", "sz": 10, "col": TS},
        {"text": "Att. v1:   42.0 %", "sz": 10, "col": ORG},
        {"text": "Att. v2:   ver hist.", "sz": 10, "col": CYN},
        {"text": " ", "sz": 3},
        {"text": "Acuerdo B↔v1 (clip a clip)", "sz": 9, "col": TM},
        {"text": "37.9 %  (37.5 % no_ident.)", "sz": 10, "col": TS},
    ], 9.25, 1.2, 3.45, 5.4)
    bx(sl, 0.5, 6.82, 8.5, 0.42, fill=PNL2)
    tx(sl, ("Lectura responsable: sin etiquetas de expertos, no se puede calcular accuracy en campo. "
            "La confianza más baja de v2 en audios OOD es DESEABLE: label smoothing + class weights reducen overconfidence, "
            "exactamente el objetivo de diseño para despliegue seguro."),
       0.65, 6.86, 8.2, 0.35, 9, TS)
    return sl


def s21(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "17. Limitaciones honestas",
        "Una evaluación científica rigorosa incluye sus propias restricciones")
    ft(sl, prs, 21)
    lims = [
        (ORG, "Desbalance parcialmente resuelto",
         "Los class weights reducen el gap de F1 entre clases, pero Celeus grammicus (F1=0.19) "
         "y Rupornis magnirostris (F1=0.25) siguen por debajo de 0.5. El desbalance en train "
         "(12:1 Trogon vs Chordeiles) supera lo que los pesos pueden compensar sin más datos."),
        (ORG, "Campo real ≠ dataset curado",
         "Los audios de Puerto Nariño tienen polución de ruido, superposición de cantos y "
         "especies no presentes en el dataset. El modelo fue entrenado sobre clips curados de "
         "Xeno-canto (grabaciones de expertos, condiciones ideales)."),
        (RED, "Salida de etiqueta única",
         "La capa softmax fuerza a elegir una clase. En audios con múltiples especies vocalizando "
         "simultáneamente, el modelo debe escoger una (o decir no_identificado). "
         "La solución es migrar a sigmoid multi-label."),
        (RED, "Detección antes de clasificación",
         "El pipeline asume que hay vocalización en el clip. Sin un VAD (Voice Activity Detector) "
         "bioacústico previo, los clips de silencio se clasifican en alguna clase, "
         "inflando la tasa de falsos positivos."),
        (TM,  "Umbral de confianza sin calibrar formalmente",
         "El umbral θ=0.6 para 'no_identificado' se eligió de forma exploratoria. "
         "Una calibración rigurosa requeriría un val set independiente con anotaciones de experto "
         "para optimizar precision-recall en cada especie."),
    ]
    for i, (col, title, desc) in enumerate(lims):
        row, col_idx = divmod(i, 3)
        if i < 3:
            x0 = 0.5 + col_idx * 4.2; y0 = 1.1
        else:
            x0 = 0.5 + (i-3)*6.25 + 0.25; y0 = 3.7
        w = 4.0 if i < 3 else 5.8
        bx(sl, x0, y0, w, 2.35, fill=PNL, border=col, bw=0.5)
        tx(sl, title, x0+0.12, y0+0.1, w-0.25, 0.32, 9.5, col, bold=True)
        tx(sl, desc,  x0+0.12, y0+0.46, w-0.25, 1.8,  8.5, TS)
    return sl


def s22(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "18. Trabajo futuro",
        "Ruta hacia un sistema de monitoreo bioacústico operativo")
    ft(sl, prs, 22)
    items = [
        (CYN, "Escalar a 50+ especies",
         "Integrar BirdCLEF 2024 + Xeno-canto Colombia. Modelo actual: 10 sp. amazónicas. "
         "Siguiente paso: cubrir los ~300 sp. del Trapecio Amazónico."),
        (GRN, "Calibración post-hoc con temperature scaling",
         "Aplicar temperature scaling (Guo et al. 2017) sobre v2 para mejorar ECE "
         "sin cambiar la arquitectura. Requiere val set independiente anotado."),
        (GLD, "Transfer learning desde BirdNET o PANNs",
         "Usar embeddings preentrenados como features y fine-tunear las últimas capas. "
         "Hipótesis: con datos muy escasos (< 20 grabaciones), el preentrenamiento ayuda."),
        (ORG, "Multi-label (sigmoid)",
         "Migrar de softmax a sigmoid para manejar clips con múltiples "
         "especies vocalizando simultáneamente. Requiere re-anotar el dataset."),
        (RED, "Despliegue en dispositivos edge",
         "Cuantizar el modelo (INT8) para correr en Raspberry Pi + micrófono. "
         "Objetivo: monitoreo continuo y autónomo en campo sin conectividad."),
    ]
    for i, (col, title, desc) in enumerate(items):
        y0 = 1.1 + i * 1.15
        bx(sl, 0.5, y0, 12.33, 1.05, fill=PNL, border=col, bw=0.5, rnd=True)
        tx(sl, title, 0.65, y0+0.08, 3.5, 0.35, 10, col, bold=True)
        tx(sl, desc,  4.3, y0+0.08, 8.3, 0.85, 9.5, TS)
    return sl


def s23(prs):
    sl = add_sl(prs); set_bg_jungle(sl, prs)
    hdr(sl, prs, "Conclusión",
        "SelvaSonic-ML: una ruta completa de ML aplicado, trazable y reproducible")
    # Columna izquierda — conclusión
    bx(sl, 0.5, 1.08, 6.5, 5.85, fill=PNL)
    tx_lines(sl, [
        {"text": "Mensaje principal", "sz": 11, "col": CYN, "bold": True},
        {"text": " ", "sz": 3},
        {"text": ("SelvaSonic-ML demuestra una ruta completa de Machine Learning aplicado: "
                  "descarga estructurada de datos → segmentación → representación Mel → "
                  "CNN baseline → +attention → +balanceo → evaluación multi-métrica → "
                  "prueba en campo."),
         "sz": 10, "col": TP},
        {"text": " ", "sz": 5},
        {"text": "Resultado cuantitativo", "sz": 11, "col": GLD, "bold": True},
        {"text": " ", "sz": 3},
        {"text": ("El modelo final (Attention v2) alcanza 0.7036 de test accuracy "
                  "sobre 11 clases con desbalance severo (12:1), sin preentrenamiento."),
         "sz": 10, "col": TP},
        {"text": " ", "sz": 5},
        {"text": "Mensaje metodológico", "sz": 11, "col": CYN, "bold": True},
        {"text": " ", "sz": 3},
        {"text": ("La mejora de un modelo no se defiende solo con accuracy: se defiende "
                  "con controles de fuga, métricas balanceadas, análisis de calibración, "
                  "embeddings interpretables y pruebas bajo cambio de distribución."),
         "sz": 10, "col": TS},
        {"text": " ", "sz": 5},
        {"text": ("Repositorio: github.com/LauraRuizA09/SelvaSonic-ML"),
         "sz": 9, "col": TM},
    ], 0.65, 1.18, 6.2, 5.65)
    # Columna derecha — referencias
    bx(sl, 7.15, 1.08, 5.68, 5.85, fill=PNL2)
    tx(sl, "Referencias", 7.3, 1.15, 5.4, 0.3, 10, CYN, bold=True)
    refs = [
        ("Gong et al. (2021).", "AST: Audio Spectrogram Transformer. Interspeech 2021."),
        ("Kong et al. (2020).", "PANNs: Large-scale pretrained audio neural networks. IEEE TASLP 28."),
        ("Piczak (2015).", "ESC: Dataset for environmental sound classification. ACM MM."),
        ("Szegedy et al. (2016).", "Rethinking Inception. CVPR 2016. [label smoothing]"),
        ("Vaswani et al. (2017).", "Attention is all you need. NeurIPS 30. [MHSA]"),
        ("Vellinga & Planqué (2015).", "The Xeno-canto collection. CLEF Working Notes."),
        ("Guo et al. (2017).", "On calibration of modern neural networks. ICML 34. [ECE]"),
        ("Cao et al. (2019).", "Learning imbalanced datasets with label-distribution-aware margin loss. NeurIPS."),
    ]
    y0_r = 1.55
    for auth, title in refs:
        y0_r_item = y0_r
        tx(sl, auth,  7.3,  y0_r_item,       5.4, 0.22, 7.5, GLD, bold=True)
        tx(sl, title, 7.3,  y0_r_item+0.2,   5.4, 0.22, 7.5, TS)
        y0_r += 0.48
    # Footer especial
    tx(sl, ("Laura Ruiz Arango & Jose Aldair Molina Méndez · Prof. Alcides Montoya · "
            "UNAL Medellín · Junio 2026"),
       0.5, 7.15, 12.33, 0.28, 8, TM, al=PP_ALIGN.CENTER)
    return sl


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    print("=== SelvaSonic-ML Presentation Builder ===")

    # 1. Generar imágenes faltantes
    print("\n[1] Generando imágenes faltantes...")
    if not IMGS["arch_baseline"].exists() or not IMGS["arch_attention"].exists():
        gen_arch_diagrams()
    else:
        print("  [ok] arch_baseline.png / arch_attention.png ya existen")
    if not JUNGLE_BG_COMPOSED.exists():
        gen_jungle_bg()
    else:
        print("  [ok] jungle_bg.png ya existe")
    if not IMGS["dist_clases"].exists():
        gen_dist_clases()
    else:
        print("  [ok] dist_clases.png ya existe")
    if not IMGS["mel_ejemplo"].exists():
        gen_mel()
    else:
        print("  [ok] mel_ejemplo.png ya existe")

    # 2. Construir presentación
    print("\n[2] Construyendo presentación (23 slides)...")
    prs = new_prs()

    builders = [
        s01, s02, s03, s04, s05, s06, s07, s08, s09, s10,
        s11, s12, s13, s14, s15, s16, s17, s18, s19, s20,
        s21, s22, s23,
    ]
    for i, fn in enumerate(builders, 1):
        fn(prs)
        print(f"  slide {i:02d} / 23  ({fn.__name__})")

    out = OUT_DIR / "SelvaSonic_ML_presentacion_final.pptx"
    prs.save(str(out))
    print(f"\n[OK] Guardado: {out}")
    print(f"     Slides: {len(prs.slides)}")

    # 3. Reporte de imágenes
    print("\n[3] Estado de imágenes embebidas:")
    for key, path in IMGS.items():
        status = "ok" if Path(path).exists() else "FALTANTE"
        print(f"  [{status:8s}] {key}: {Path(path).name}")


if __name__ == "__main__":
    main()
